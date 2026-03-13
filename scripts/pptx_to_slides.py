#!/usr/bin/env python3
"""
pptx_to_html.py — PPTX 转 HTML Slides
Visual Cognition Skill · 爱思考的伊伊子

用法：
  python3 pptx_to_html.py input.pptx [output.html] [--theme 1]

依赖：
  pip install python-pptx Pillow
"""

import sys
import json
import base64
import io
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from PIL import Image
except ImportError:
    print("请先安装依赖：pip install python-pptx Pillow")
    sys.exit(1)


def rgb_to_hex(rgb):
    if rgb is None:
        return None
    return '#{:02x}{:02x}{:02x}'.format(rgb.r, rgb.g, rgb.b)


def image_to_base64(image_bytes, format='PNG'):
    img = Image.open(io.BytesIO(image_bytes))
    # 限制最大尺寸
    img.thumbnail((1920, 1080), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format=format)
    return base64.b64encode(buf.getvalue()).decode('utf-8')


def extract_slide(slide, slide_idx):
    """提取单个 slide 的结构化内容"""
    result = {
        'idx':     slide_idx,
        'title':   None,
        'texts':   [],   # [{text, size, bold, align, color}]
        'images':  [],   # [{base64, width_pct, height_pct, left_pct, top_pct}]
        'bg_color':None,
        'notes':   None,
    }

    # 背景色
    try:
        bg = slide.background.fill
        if bg.fore_color:
            result['bg_color'] = rgb_to_hex(bg.fore_color.rgb)
    except:
        pass

    # 演讲者备注
    try:
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            result['notes'] = notes_tf.text.strip()
    except:
        pass

    # 形状
    slide_w = slide.shapes._spTree.getparent().getparent().slide_width
    slide_h = slide.shapes._spTree.getparent().getparent().slide_height

    for shape in slide.shapes:
        # 位置（百分比）
        l = shape.left   / slide_w * 100 if shape.left   else 0
        t = shape.top    / slide_h * 100 if shape.top    else 0
        w = shape.width  / slide_w * 100 if shape.width  else 0
        h = shape.height / slide_h * 100 if shape.height else 0

        # 文字
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # 字号和样式（取第一个run）
                font_size = 24
                bold      = False
                color     = None
                align     = 'left'

                if para.runs:
                    run = para.runs[0]
                    if run.font.size:
                        font_size = int(run.font.size / 12700)  # EMU → pt
                    bold  = run.font.bold or False
                    color = rgb_to_hex(run.font.color.rgb) if run.font.color and run.font.color.type else None

                if para.alignment == PP_ALIGN.CENTER:
                    align = 'center'
                elif para.alignment == PP_ALIGN.RIGHT:
                    align = 'right'

                entry = {
                    'text':      text,
                    'font_size': font_size,
                    'bold':      bold,
                    'color':     color,
                    'align':     align,
                    'left_pct':  round(l, 1),
                    'top_pct':   round(t, 1),
                    'w_pct':     round(w, 1),
                }

                # 判断是否是标题
                if shape.shape_type == 13 or (font_size > 28 and not result['title']):
                    result['title'] = text
                else:
                    result['texts'].append(entry)

        # 图片
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                img_bytes = shape.image.blob
                b64 = image_to_base64(img_bytes)
                result['images'].append({
                    'base64':   b64,
                    'format':   shape.image.content_type.split('/')[-1].upper(),
                    'left_pct': round(l, 1),
                    'top_pct':  round(t, 1),
                    'w_pct':    round(w, 1),
                    'h_pct':    round(h, 1),
                })
            except Exception as e:
                print(f"  图片提取失败（Slide {slide_idx+1}）: {e}")

    return result


def slides_to_html(slides_data, theme=1, title="Presentation"):
    """将提取的 slide 数据转换为 HTML"""

    slides_html = []
    for i, s in enumerate(slides_data):
        # 主标题
        title_html = f'<div class="ex-title editable">{s["title"]}</div>' if s['title'] else ''

        # 文字块
        texts_html = ''
        for t in s['texts']:
            size_scale = min(max(t['font_size'] / 24, 0.8), 2.5)
            style = f'font-size:calc(var(--size-body) * {size_scale:.1f});'
            if t['bold']:     style += 'font-weight:900;'
            if t['color']:    style += f'color:{t["color"]};'
            if t['align'] != 'left': style += f'text-align:{t["align"]};'
            texts_html += f'<p class="ex-text editable" style="{style}">{t["text"]}</p>\n'

        # 图片
        imgs_html = ''
        for img in s['images']:
            fmt = img['format'].lower()
            if fmt == 'jpeg': fmt = 'jpg'
            imgs_html += f'<img src="data:image/{fmt};base64,{img["base64"]}" class="ex-img" alt="slide image">\n'

        # 备注（作为 data 属性，不显示）
        notes_attr = f' data-notes="{s["notes"]}"' if s['notes'] else ''

        slide_html = f'''
    <!-- Slide {i+1} -->
    <section class="slide{'  active' if i==0 else ''}" id="s{i+1}"{notes_attr}>
      <div class="ex-content">
        {title_html}
        {texts_html}
        {imgs_html}
      </div>
      <div class="sn">{i+1:02d} / {len(slides_data):02d}</div>
    </section>'''
        slides_html.append(slide_html)

    all_slides = '\n'.join(slides_html)

    return f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{title}</title>
<link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@400;900&family=Space+Mono&display=swap" rel="stylesheet">
<style>
/* ── 由 Visual Cognition Skill 从 PPTX 转换 ── */
:root {{
  --bg:#f5f0e8; --ink:#1a1410;
  --size-body:clamp(14px,1.35vw,28px);
  --size-h1:clamp(24px,2.7vw,56px);
  --size-display:clamp(36px,4vw,88px);
}}
*{{box-sizing:border-box;margin:0;padding:0;}}
html,body{{width:100%;height:100%;overflow:hidden;background:#111;}}
.deck-wrap{{position:fixed;inset:0;display:flex;align-items:center;justify-content:center;}}
.deck{{width:1920px;height:1080px;position:relative;transform-origin:top left;}}
.slide{{
  position:absolute;inset:0;background:var(--bg);
  display:none;flex-direction:column;align-items:center;justify-content:center;
  overflow:hidden;
}}
.slide::before{{
  content:'';position:absolute;inset:0;
  background-image:radial-gradient(circle,rgba(26,20,16,.1) 1.5px,transparent 1.5px);
  background-size:40px 40px;pointer-events:none;
}}
.slide.active{{display:flex;}}
.ex-content{{
  position:relative;z-index:1;
  width:90%;max-height:90%;
  display:flex;flex-direction:column;gap:24px;
}}
.ex-title{{
  font-family:'Noto Sans SC',sans-serif;font-weight:900;
  font-size:var(--size-display);line-height:1.15;
}}
.ex-text{{
  font-family:'Noto Sans SC',sans-serif;
  font-size:var(--size-body);line-height:1.6;
}}
.ex-img{{max-width:80%;max-height:400px;object-fit:contain;align-self:center;}}
.sn{{
  position:absolute;bottom:24px;right:40px;
  font-family:'Space Mono',monospace;font-size:14px;opacity:.25;
}}
#prog{{position:absolute;bottom:0;left:0;height:4px;background:var(--ink);transition:width .3s;z-index:200;}}

/* 编辑模式 */
[contenteditable=true]{{outline:2px dashed #e8a020!important;min-width:20px;min-height:20px;}}
#editBar{{
  position:fixed;bottom:20px;left:50%;transform:translateX(-50%);
  background:#1a1410;color:#f5f0e8;
  padding:10px 20px;border-radius:999px;z-index:999;
  display:none;gap:12px;
}}
#editBar button{{
  background:none;border:1px solid rgba(255,255,255,.3);
  color:inherit;padding:6px 16px;border-radius:999px;cursor:pointer;font-size:14px;
}}
</style>
</head>
<body>
<div class="deck-wrap">
<div class="deck" id="deck" data-w="1920" data-h="1080">
{all_slides}
<div id="prog"></div>
</div>
</div>

<!-- 编辑工具栏（按 E 键打开）-->
<div id="editBar">
  <span style="font-size:14px;opacity:.5">编辑模式</span>
  <button onclick="toggleEdit()">✏️ 切换</button>
  <button onclick="exportHTML()">⬇️ 导出</button>
  <button onclick="document.getElementById('editBar').style.display='none'">✕</button>
</div>

<script>
// 缩放
function scaleCanvas() {{
  const d=document.getElementById('deck');
  const W=+d.dataset.w||1920, H=+d.dataset.h||1080;
  const s=Math.min(innerWidth/W, innerHeight/H);
  d.style.cssText=`width:${{W}}px;height:${{H}}px;transform:scale(${{s}});transform-origin:top left;margin-left:${{(innerWidth-W*s)/2}}px;margin-top:${{(innerHeight-H*s)/2}}px;`;
}}
window.addEventListener('resize',scaleCanvas); scaleCanvas();

// 导航
const slides=[...document.querySelectorAll('.slide')];
const prog=document.getElementById('prog');
let cur=0;
function goto(n){{
  slides[cur].classList.remove('active');
  cur=Math.max(0,Math.min(n,slides.length-1));
  slides[cur].classList.add('active');
  prog.style.width=((cur+1)/slides.length*100)+'%';
}}
document.addEventListener('keydown',e=>{{
  if(e.key==='ArrowRight'||e.key===' '){{e.preventDefault();goto(cur+1);}}
  if(e.key==='ArrowLeft'){{e.preventDefault();goto(cur-1);}}
  if(e.key==='e'){{document.getElementById('editBar').style.display='flex';}}
}});
document.addEventListener('click',e=>{{
  if(e.target.closest('#editBar,button,a'))return;
  goto(cur+1);
}});
let tx=0;
document.addEventListener('touchstart',e=>{{tx=e.touches[0].clientX;}});
document.addEventListener('touchend',e=>{{
  const dx=e.changedTouches[0].clientX-tx;
  if(Math.abs(dx)>50){{dx<0?goto(cur+1):goto(cur-1);}}
}});

// 编辑
let editMode=false;
function toggleEdit(){{
  editMode=!editMode;
  document.querySelectorAll('.editable').forEach(el=>{{
    el.contentEditable=editMode?'true':'false';
  }});
}}
function exportHTML(){{
  const html='<!DOCTYPE html>\\n'+document.documentElement.outerHTML;
  const a=document.createElement('a');
  a.href=URL.createObjectURL(new Blob([html],{{type:'text/html'}}));
  a.download='slides-exported.html'; a.click();
}}
</script>
</body>
</html>'''


def main():
    args = sys.argv[1:]
    if not args:
        print("用法：python3 pptx_to_html.py input.pptx [output.html] [--theme N]")
        sys.exit(1)

    input_path  = Path(args[0])
    output_path = Path(args[1]) if len(args) > 1 and not args[1].startswith('--') else input_path.with_suffix('.html')
    theme       = 1

    for i, a in enumerate(args):
        if a == '--theme' and i+1 < len(args):
            theme = int(args[i+1])

    print(f"📂 读取：{input_path}")
    prs = Presentation(str(input_path))

    print(f"📊 共 {len(prs.slides)} 张 slides，开始提取…")
    slides_data = []
    for i, slide in enumerate(prs.slides):
        print(f"  Slide {i+1}/{len(prs.slides)}")
        slides_data.append(extract_slide(slide, i))

    print("🎨 生成 HTML…")
    html = slides_to_html(slides_data, theme=theme, title=input_path.stem)

    output_path.write_text(html, encoding='utf-8')
    print(f"✅ 完成：{output_path}")
    print(f"   提示：在浏览器打开，按 E 键进入编辑模式，⬇️ 导出修改后的版本")


if __name__ == '__main__':
    main()
